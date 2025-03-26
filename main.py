import os
#import pyautogui
from pptx import Presentation
from pptx.util import Inches
from mpath import main_path

l = 1
t = 0
w = 8
h = 7.5

list = os.listdir(main_path)
list = sorted(list)
root = Presentation()
for i in list:
    img_path = main_path + '/'+ i
    black_slide = root.slide_layouts[6]
    slide = root.slides.add_slide(black_slide)

    left = Inches(l)
    top = Inches(t)
    width = Inches(w)
    height = Inches(h)

    pic = slide.shapes.add_picture(img_path, left, top, width, height)


root.save('Test.pptx')
